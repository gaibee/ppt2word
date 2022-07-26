from pptx import Presentation
from pptx_tools import utils

def slide_to_image(png_folder, ppt_path):
    '''
    you should use full paths, to make sure PowerPoint can handle the paths
    '''
    utils.save_pptx_as_png(png_folder, f'"{ppt_path}"', overwrite_folder=True)


def get_slides_note(ppt_path):
    parsed = Presentation(ppt_path)
    notes = []

    # 슬라이드 별로 순회하면서 데이터 추출
    for slide in parsed.slides:
        # silde_note = slide.notes_slide.notes_text_frame.text
        notes.append(slide.notes_slide.notes_text_frame.text)

    return notes


if __name__=='__main__':
    '''
    parsed = Presentation("data/myppt.pptx")

    # 슬라이드 별로 순회하면서 데이터 추출
    for i, slide in enumerate(parsed.slides):
        # silde_note = slide.notes_slide.notes_text_frame.text
        print(i, slide.name)
    '''
    